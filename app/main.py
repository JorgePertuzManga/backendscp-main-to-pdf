# app/main.py
from fastapi import FastAPI, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
# from pathlib import Path
from pathlib import Path # Path no importado (rompe la app al arrancar
import io

from .models import RenderRequest
from .ppt import generate_presentation
# from .pdf_converter import PDFConverter, OnlinePDFConverter
import tempfile
import subprocess
import os

TEMPLATE_PATH = str(Path(__file__).parent / "templates" / "silicon_eic_template.pptx")

app = FastAPI()

# Ajusta orígenes si quieres restringirlo
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def _build_replacements(req: RenderRequest) -> dict:
    """
    Solo usamos los placeholders que tu plantilla necesita.
    Enviamos las llaves con {{...}} para que el replace sea exacto.
    """
    repl = {
        "{{COMPANY_NAME}}": req.company_name,
    }

    # pricing_overrides: solo añadimos los que vengan
    po = req.pricing_overrides or {}
    if po.SETUP_FEE is not None:
        repl["{{SETUP_FEE}}"] = str(po.SETUP_FEE)
    if po.SHORT_FEE is not None:
        repl["{{SHORT_FEE}}"] = str(po.SHORT_FEE)
    if po.FULL_FEE is not None:
        repl["{{FULL_FEE}}"] = str(po.FULL_FEE)
    if po.GRANT_FEE is not None:
        repl["{{GRANT_FEE}}"] = po.GRANT_FEE
    if po.EQUITY_FEE is not None:
        repl["{{EQUITY_FEE}}"] = po.EQUITY_FEE

    return repl


@app.post("/render")
async def render_presentation(request: RenderRequest):
    try:
        replacements = _build_replacements(request)
        buf = generate_presentation(
            template_path=TEMPLATE_PATH,
            replacements=replacements,
            slide_toggles=request.slide_toggles or {},
        )
        filename = f"proposal_{request.company_name.replace(' ', '_')}.pptx"
        return StreamingResponse(
            io.BytesIO(buf.getvalue()),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
    

@app.post("/render_pdf")
async def render_presentation_pdf(request: RenderRequest):
    """
    Nuevo endpoint que devuelve pdf en vez de pptx
    """
    try:
        if not os.path.exists(TEMPLATE_PATH):
            raise HTTPException(status_code=500, detail="Template file not found.")
        """
        overrides = requests.pricing_overrides
        replacements = {
            "{COMPANY_NAME}": request.company_name,
            "{SETUP_FEE}": f"{overrides.SETUP_FEE: ,.0f}$".replace(",",""),
            "{SHORT_FEE}": f"{overrides.SHORT_FEE: ,.0f}$".replace(",",""),
            "{FULL_FEE}": f"{overrides.FULL_FEE: ,.0f}$".replace(",",""),
            "{GRANT_FEE}": overrides.GRANT_FEE,
            "{EQUITY_FEE}": overrides.EQUITY_FEE,
            "{DATE}": request.proposal_date.strftime("%B %d, %Y"),
        }
        """
        # 1) Reutilizar los mismos placeholders que usa /render
        replacements = _build_replacements(request)

        # 2) Agregar DATE solo si la plantilla lo usa
        if request.proposal_date:
            replacements["{{DATE}}"] = request.proposal_date.strftime("%B %d, %Y")
        else:
            replacements["{{DATE}}"] = ""
            
        pptx_stream = generate_presentation(
            template_path=TEMPLATE_PATH,
            replacements=replacements,
            slide_toggles=request.slide_toggles or {} # Agregar  or {}
        )

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as pptx_temp:
            pptx_temp.write(pptx_stream.getvalue())
            pptx_temp.flush()
            pptx_path = pptx_temp.name

        pdf_path = pptx_path.replace(".pptx", ".pdf")

        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to",
            "pdf",
            pptx_path,
            "--outdir",
            os.path.dirname(pptx_path)
        ], check=True)

        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()
        
        os.remove(pptx_path)

        os.remove(pdf_path)

        safe_company = "".join(c if c.isalnum() or c in " _-" else "_" for c in request.company_name)
        filename = f"proposal_eic_template_{safe_company}.pdf"

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'},
        )
    
    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"error converting to libreoffice: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"error general: {str(e)}")
    
    
@app.post("/render_pdf_custom")
async def render_pdf_custom(request: RenderRequest, remove_slides: list[int] = []):
    """
    Genera un PDF desde la plantilla PPTX permitiendo al usuario eliminar
    páginas específicas por número de slide (1-based index).
    
    Ejemplo:
    remove_slides = [2, 5, 7]
    """

    try:
        if not os.path.exists(TEMPLATE_PATH):
            raise HTTPException(status_code=500, detail="Template file not found.")

        # 1) Construimos los replacements igual que /render y /render_pdf
        replacements = _build_replacements(request)

        # Fecha opcional
        if request.proposal_date:
            replacements["{{DATE}}"] = request.proposal_date.strftime("%B %d, %Y")
        else:
            replacements["{{DATE}}"] = ""

        # 2) Convertir slides enviados por el usuario (1-based → 0-based)
        #    Si pide eliminar slide 2 → eliminamos índice 1
        remove_idx = sorted([s - 1 for s in remove_slides if s > 0])

        # 3) Generar primero la presentación PPTX completa en memoria
        pptx_stream = generate_presentation(
            template_path=TEMPLATE_PATH,
            replacements=replacements,
            slide_toggles=request.slide_toggles or {}
        )

        # 4) Guardarla temporalmente para poder manipularla
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as temp_ppt:
            temp_ppt.write(pptx_stream.getvalue())
            temp_ppt.flush()
            ppt_path = temp_ppt.name

        # 5) Abrir la presentación para eliminar páginas específicas
        from pptx import Presentation
        prs = Presentation(ppt_path)

        # Importante: borrar desde atrás hacia adelante
        for idx in reversed(remove_idx):
            if 0 <= idx < len(prs.slides._sldIdLst):
                rId = prs.slides._sldIdLst[idx].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]

        # 6) Guardar el PPTX ya modificado (sin slides)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as final_ppt:
            prs.save(final_ppt.name)
            final_ppt_path = final_ppt.name

        # 7) Convertir a PDF con LibreOffice
        pdf_path = final_ppt_path.replace(".pptx", ".pdf")

        subprocess.run([
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            final_ppt_path,
            "--outdir", os.path.dirname(final_ppt_path),
        ], check=True)

        # Leer PDF
        with open(pdf_path, "rb") as f:
            pdf_bytes = f.read()

        # Limpiar archivos
        os.remove(ppt_path)
        os.remove(final_ppt_path)
        os.remove(pdf_path)

        # Nombre final
        safe_company = "".join(c if c.isalnum() or c in " _-" else "_" for c in request.company_name)
        filename = f"proposal_custom_{safe_company}.pdf"

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename=\"{filename}\"'}
        )

    except subprocess.CalledProcessError as e:
        raise HTTPException(status_code=500, detail=f"LibreOffice conversion error: {str(e)}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"General error: {str(e)}")

