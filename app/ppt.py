from pptx import Presentation
from io import BytesIO
from typing import Dict
import re


def replace_placeholders_in_text(text_frame, replacements: Dict[str, str]) -> None:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            for key, value in replacements.items():
                if key in run.text:
                    run.text = run.text.replace(key, value)


def get_slide_tags(notes_text: str) -> list:
    # Solo letras, números y _ ; evita el rango A-z que incluye símbolos.
    pattern = r"\[\[tag:([A-Za-z0-9_]+)\]\]"
    return re.findall(pattern, notes_text or "")


def generate_presentation(
    template_path: str, replacements: Dict[str, str], slide_toggles: Dict[str, bool]
) -> BytesIO:
    prs = Presentation(template_path)

    # 1) Marcar slides a eliminar por tags en NOTAS
    slides_to_remove = []
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            tags = get_slide_tags(slide.notes_slide.notes_text_frame.text)
            should_remove = any(not slide_toggles.get(tag, True) for tag in tags)
            if should_remove:
                slides_to_remove.append(i)

    # 2) Eliminar (de atrás hacia delante)
    for i in reversed(slides_to_remove):
        rId = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[i]

    # 3) Reemplazar placeholders en todas las shapes con texto + también en notas
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                replace_placeholders_in_text(shape.text_frame, replacements)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            replace_placeholders_in_text(
                slide.notes_slide.notes_text_frame, replacements
            )

    # 4) Guardar en memoria
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
